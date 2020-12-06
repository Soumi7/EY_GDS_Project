import os
import re
import PyPDF2
import csv
import csv
import nltk
import sklearn
import numpy as np
import pandas as pd
from sklearn.utils import shuffle
from io import StringIO

from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfdocument import PDFDocument
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.pdfpage import PDFPage
from pdfminer.pdfparser import PDFParser

from nltk.tokenize import sent_tokenize
nltk.download('punkt')

from pptx import Presentation

import docx2txt
import docx

csv_data=[]
file_loc=[]
sentence =[]
label=[]
intent=[]
file_names=[]
file_name=[]

dir="/home/lohith/Desktop/EY Hackathon/EY_GDS_Project/data_to_preprocess"
def list_files(dir):
    r = []
    for root, dirs, files in os.walk(dir):
        for name in files:
            r.append(os.path.join(root, name))
    return r
for i in list_files(dir):
    file_names.append(i)

# print(file_location)

for filename in file_names:
    # print(filename)
    if filename.endswith('.pdf'):
        output_string = StringIO()
        with open(filename, 'rb') as in_file:
            parser = PDFParser(in_file)
            doc = PDFDocument(parser)
            rsrcmgr = PDFResourceManager()
            device = TextConverter(rsrcmgr, output_string, laparams = LAParams())
            interpreter = PDFPageInterpreter(rsrcmgr, device)
            for page in PDFPage.create_pages(doc):
                interpreter.process_page(page)

        text=output_string.getvalue()
        data = sent_tokenize(text)

        for line in data:
            res=re.sub('\s+',' ',line)


            line=str(res)
            line = re.sub(r'[^\w\s]', '', line)
            line = re.sub(r'\b(?!(\D\S*|[12][0-9]{3})\b)\S+\b', '', line)
            line = line.strip()
            file_loc.append(filename)

            try:
                type(int(line)) != int
 
            except ValueError:
                if len(line) > 10:
                    sentence.append(line)

            intent.append(filename.split("/")[7])
            file_name.append(filename.split("/")[-1])

    # if filename.endswith('.pptx'):
    if filename.endswith('.pptx') :
        print(filename)
        output_string = StringIO()
        with open(filename, 'rb') as in_file:
            prs = Presentation(in_file)
            content = ""
            slideCount = 0
            fullText = []
            for slide in prs.slides:
                slideCount += 1
                for shape in slide.shapes:
                    if (shape.has_text_frame):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                fullText.append(run.text)
            fullText='\n'.join(fullText)

        text=str(fullText)
        data = sent_tokenize(text)
        for line in data:
            res=re.sub('\s+',' ',line)
            line=str(res)

            line = re.sub(r'[^\w\s]', '', line)
            line = re.sub(r"\d+", "", line)
            # print(line)
            if len(line) != 0:
                file_loc.append(filename)
                try:
                    type(int(line)) != int

                except ValueError:
                    sentence.append(line)

                intent.append(filename.split("/")[-2])
                file_name.append(filename.split("/")[-1])

    if filename.endswith('.docx'):
            output_string = StringIO()
            with open(filename, 'rb') as in_file:
                doc = docx.Document(in_file)
                fullText = []
                for para in doc.paragraphs:
                    fullText.append(para.text)
                fullText='\n'.join(fullText)

            text=str(fullText)
            data = sent_tokenize(text)
            for line in data:
                res=re.sub('\s+',' ',line)
                line=str(res)

                line = re.sub(r'[^\w\s]', '', line)
                line = re.sub(r"\d+", "", line)
                # print(line)
                if len(line) != 0:
                    file_loc.append(filename)
                    try:
                        type(int(line)) != int
    
                    except ValueError:
                        sentence.append(line)

                    
                    intent.append(filename.split("/")[-2])
                    file_name.append(filename.split("/")[-1])

df = pd.DataFrame(list(zip(file_loc, file_name, sentence , intent)) , columns=["File Location", "File Name", "Sentence" , "Intent"])
df.to_csv('folders.csv',encoding='utf-8-sig', index=False) 

df_compressed =  df[["Sentence","Intent"]]

df_compressed =  df_compressed.rename({'Sentence':'text', "Intent":"intent"}, axis=1) 

df_compressed = shuffle(df_compressed)

print(df["Intent"].value_counts())

size =  len(df_compressed)

train_df = df_compressed[:int(np.round(size*.6, 0))]
test_df = df_compressed[int(np.round(size*.6, 0)) : int(np.round(size*.6, 0))+int(np.round(size*.2, 0))]
valid_df = df_compressed[int(np.round(size*.6, 0))+int(np.round(size*.2, 0)):]


train_df.to_csv('train.csv',encoding='utf-8-sig', index=False) 
test_df.to_csv('test.csv',encoding='utf-8-sig', index=False) 
valid_df.to_csv('valid.csv',encoding='utf-8-sig', index=False) 
