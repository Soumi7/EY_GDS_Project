

from pptx import Presentation


prs = Presentation("/home/soumi/Downloads/EYintentRepository/Repository 1/Item 4/azure-250215.pptx")
print("----------------------")
'''
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(shape.text)
'''
content = ""
slideCount = 0
for slide in prs.slides:
    slideCount += 1
    for shape in slide.shapes:
        if (shape.has_text_frame):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    print(run.text)