from pptx import Presentation
import os
import PyPDF2
import data_func
import csv
import pandas as pd
from io import StringIO
import re
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfdocument import PDFDocument
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.pdfpage import PDFPage
from pdfminer.pdfparser import PDFParser
from nltk.tokenize import sent_tokenize
import csv
import nltk
# nltk.download('punkt')
import sklearn
import numpy as np
from sklearn.utils import shuffle
import docx2txt
import docx


csv_data=[]
file_loc=[]
sentence =[]
label=[]
intent=[]
file_names=[]
file_name=[]

dir="/home/lohith/Desktop/EY Hackathon/EY_GDS_Project/data_to_preprocess"
def list_files(dir):
	r = []
	for root, dirs, files in os.walk(dir):
		for name in files:
			r.append(os.path.join(root, name))
	return r
for i in list_files(dir):
	file_names.append(i)

# print(file_location)

for filename in file_names:
	# print(filename)
	if filename.endswith('.pptx') :
		print(filename)
		output_string = StringIO()
		with open(filename, 'rb') as in_file:
			prs = Presentation(in_file)
			content = ""
			slideCount = 0
			fullText = []
			for slide in prs.slides:
				slideCount += 1
				for shape in slide.shapes:
					if (shape.has_text_frame):
						for paragraph in shape.text_frame.paragraphs:
							for run in paragraph.runs:
								fullText.append(run.text)
			fullText='\n'.join(fullText)

		text=str(fullText)
		data = sent_tokenize(text)
		for line in data:
			res=re.sub('\s+',' ',line)
			line=str(res)

			line = re.sub(r'[^\w\s]', '', line)
			line = re.sub(r"\d+", "", line)
			# print(line)
			if len(line) != 0:
				file_loc.append(filename)
				try:
					type(int(line)) != int
 
				except ValueError:
					sentence.append(line)

				
				label.append(0)
				intent.append(filename.split("/")[-2])
				file_name.append(filename.split("/")[-1])


df = pd.DataFrame(list(zip(file_loc, file_name, sentence , label, intent)) , columns=["File Location", "File Name", "Sentence" , "Label", "Intent"])
df.to_csv('PPtx.csv',encoding='utf-8-sig', index=False) 

# df = pd.read_csv("folders.csv", encoding="utf-8-sig")

df_compressed =  df[["Sentence","Intent"]]

df_compressed =  df_compressed.rename({'Sentence':'text', "Intent":"intent"}, axis=1) 

df_compressed = shuffle(df_compressed)



size =  len(df_compressed)

train_df = df_compressed[:int(np.round(size*.6, 0))]
test_df = df_compressed[int(np.round(size*.6, 0)) : int(np.round(size*.6, 0))+int(np.round(size*.2, 0))]
valid_df = df_compressed[int(np.round(size*.6, 0))+int(np.round(size*.2, 0)):]


train_df.to_csv('train_ppt.csv',encoding='utf-8-sig', index=False) 
test_df.to_csv('test_ppt.csv',encoding='utf-8-sig', index=False) 
valid_df.to_csv('valid_ppt.csv',encoding='utf-8-sig', index=False) 

